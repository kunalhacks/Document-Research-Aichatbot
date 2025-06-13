import os
import io
import mimetypes
import pytesseract
from typing import List, Optional, Tuple, Union
from pathlib import Path
from PIL import Image
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import cv2
import numpy as np
import pytesseract
from concurrent.futures import ThreadPoolExecutor
import logging

try:
    import magic
except Exception:
    magic = None

from app.models.document import Document, DocumentChunk, DocumentMetadata, DocumentType

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Handles processing of different document types and extracts text content."""
    
    def __init__(self, tesseract_path: Optional[str] = None):
        """Initialize the document processor.
        
        Args:
            tesseract_path: Path to Tesseract OCR executable if not in system PATH
        """
        if tesseract_path:
            pytesseract.pytesseract.tesseract_cmd = tesseract_path
        
        self.supported_mime_types = {
            'application/pdf': self._process_pdf,
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': self._process_docx,
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': self._process_pptx,
            'text/plain': self._process_text,
            'image/jpeg': self._process_image,
            'image/png': self._process_image,
            'image/tiff': self._process_image,
        }
    
    def get_document_type(self, file_path: Union[str, Path]) -> Optional[Tuple[str, str]]:
        """Determine the MIME type and file extension of a file."""
        try:
            if magic is not None:
                mime = magic.Magic(mime=True)
                mime_type = mime.from_file(str(file_path))
            else:
                mime_type, _ = mimetypes.guess_type(str(file_path))
            ext = Path(file_path).suffix.lower()
            return mime_type or 'application/octet-stream', ext[1:] if ext else None
        except Exception as e:
            logger.error(f"Error determining file type: {e}")
            return None, None
    
    def process_document(self, file_path: Union[str, Path]) -> Optional[Document]:
        """Process a document and return a Document object with extracted text."""
        try:
            file_path = Path(file_path)
            if not file_path.exists():
                raise FileNotFoundError(f"File not found: {file_path}")
                
            mime_type, ext = self.get_document_type(file_path)
            if not mime_type:
                raise ValueError(f"Could not determine file type for {file_path}")
                
            logger.info(f"Processing {file_path} (MIME: {mime_type}, Ext: {ext})")
            
            # Get the appropriate processor function
            processor = self.supported_mime_types.get(mime_type)
            if not processor:
                raise ValueError(f"Unsupported file type: {mime_type}")
            
            # Process the document
            text_chunks = processor(file_path)
            if not text_chunks:
                raise ValueError("No text could be extracted from the document")
            
            # Create document metadata
            doc_metadata = DocumentMetadata(
                doc_id=str(file_path.stem),
                title=file_path.stem,
                doc_type=DocumentType(ext) if ext in [t.value for t in DocumentType] else DocumentType.TEXT
            )
            
            # Create document with chunks
            document = Document(metadata=doc_metadata)
            for i, (page_num, chunk_text) in enumerate(text_chunks):
                chunk = DocumentChunk(
                    chunk_id=f"{doc_metadata.doc_id}_chunk_{i}",
                    doc_id=doc_metadata.doc_id,
                    content=chunk_text,
                    page_number=page_num,
                    chunk_index=i,
                    metadata={"source": str(file_path.name), "page": str(page_num + 1)}
                )
                document.add_chunk(chunk)
            
            document.metadata.page_count = max([c.page_number for c in document.chunks], default=0) + 1
            document.metadata.word_count = sum(len(chunk.content.split()) for chunk in document.chunks)
            
            return document
            
        except Exception as e:
            logger.error(f"Error processing document {file_path}: {e}", exc_info=True)
            return None
    
    def _process_pdf(self, file_path: Union[str, Path]) -> List[Tuple[int, str]]:
        """Extract text from PDF file."""
        try:
            from pdf2image import convert_from_path
            
            # First try to extract text directly
            try:
                reader = PdfReader(file_path)
                chunks = []
                for i, page in enumerate(reader.pages):
                    text = page.extract_text()
                    if text.strip():
                        chunks.append((i, text))
                
                if chunks:
                    return chunks
            except Exception as e:
                logger.warning(f"Direct PDF text extraction failed, falling back to OCR: {e}")
            
            # Fall back to OCR if direct extraction fails
            images = convert_from_path(file_path)
            chunks = []
            for i, image in enumerate(images):
                text = pytesseract.image_to_string(image)
                if text.strip():
                    chunks.append((i, text))
            return chunks
            
        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {e}")
            return []
    
    def _process_docx(self, file_path: Union[str, Path]) -> List[Tuple[int, str]]:
        """Extract text from DOCX file."""
        try:
            doc = DocxDocument(file_path)
            full_text = []
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text)
            
            if full_text:
                return [(0, "\n".join(full_text))]
            return []
        except Exception as e:
            logger.error(f"Error processing DOCX {file_path}: {e}")
            return []
    
    def _process_pptx(self, file_path: Union[str, Path]) -> List[Tuple[int, str]]:
        """Extract text from PPTX file."""
        try:
            prs = Presentation(file_path)
            slides_text = []
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                if slide_text:
                    slides_text.append((i, "\n".join(slide_text)))
            return slides_text
        except Exception as e:
            logger.error(f"Error processing PPTX {file_path}: {e}")
            return []
    
    def _process_text(self, file_path: Union[str, Path]) -> List[Tuple[int, str]]:
        """Extract text from plain text file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
                return [(0, content)] if content.strip() else []
        except Exception as e:
            logger.error(f"Error processing text file {file_path}: {e}")
            return []
    
    def _process_image(self, file_path: Union[str, Path]) -> List[Tuple[int, str]]:
        """Extract text from image file using OCR."""
        try:
            # Read the image
            image = cv2.imread(str(file_path))
            if image is None:
                raise ValueError(f"Could not read image: {file_path}")
                
            # Convert to grayscale
            gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
            
            # Apply thresholding to preprocess the image
            gray = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)[1]
            
            # Perform OCR
            text = pytesseract.image_to_string(gray)
            
            return [(0, text)] if text.strip() else []
            
        except Exception as e:
            logger.error(f"Error processing image {file_path}: {e}")
            return []
    
    def process_directory(self, dir_path: Union[str, Path], recursive: bool = True) -> List[Document]:
        """Process all supported documents in a directory."""
        dir_path = Path(dir_path)
        if not dir_path.is_dir():
            raise NotADirectoryError(f"Directory not found: {dir_path}")
        
        documents = []
        files_to_process = []
        
        # Collect all supported files
        for ext in [".pdf", ".docx", ".pptx", ".txt", ".jpg", ".jpeg", ".png", ".tiff"]:
            pattern = "**/*" + ext if recursive else "*" + ext
            files_to_process.extend(dir_path.glob(pattern))
        
        # Process files in parallel
        with ThreadPoolExecutor() as executor:
            results = list(executor.map(self.process_document, files_to_process))
        
        # Filter out None results
        documents = [doc for doc in results if doc is not None]
        
        return documents
